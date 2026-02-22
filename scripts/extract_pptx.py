import sys
import os
import json
try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx library not found. Install it with: pip install python-pptx")
    sys.exit(1)

def extract_pptx(file_path, output_dir):
    """
    Extract all content from a PowerPoint file.
    Returns a JSON structure with slides, text, and images.
    """
    if not os.path.exists(file_path):
        print(f"Error: File not found {file_path}")
        return None

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error: Failed to open PowerPoint file: {e}")
        return None

    slides_data = []

    # Create assets directory
    assets_dir = os.path.join(output_dir, 'assets')
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            'number': slide_num + 1,
            'title': '',
            'content': [],
            'images': [],
            'notes': ''
        }

        for shape in slide.shapes:
            # Extract title and text
            if shape.has_text_frame:
                text_content = shape.text.strip()
                if not text_content:
                    continue
                
                if shape == slide.shapes.title:
                    slide_data['title'] = text_content
                else:
                    slide_data['content'].append({
                        'type': 'text',
                        'content': text_content
                    })

            # Extract images
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                    image_path = os.path.join(assets_dir, image_name)

                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)

                    slide_data['images'].append({
                        'path': f"assets/{image_name}",
                        'width': shape.width,
                        'height': shape.height
                    })
                except Exception as e:
                    print(f"Warning: Failed to extract image from slide {slide_num+1}: {e}")

        # Extract notes
        if slide.has_notes_slide:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                slide_data['notes'] = notes_frame.text
            except:
                pass

        slides_data.append(slide_data)

    return slides_data

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx.py <path_to_pptx> <output_directory>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    out_dir = sys.argv[2]

    data = extract_pptx(pptx_path, out_dir)
    if data:
        with open(os.path.join(out_dir, 'content.json'), 'w') as f:
            json.dump(data, f, indent=2)
        print(f"Successfully extracted {len(data)} slides to {out_dir}")
